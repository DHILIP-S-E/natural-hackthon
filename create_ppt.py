"""Generate AURA Hackathon Pitch Deck — 8 Slides"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color palette ──
BG_DARK = RGBColor(0x0A, 0x0A, 0x0F)
BG_CARD = RGBColor(0x16, 0x16, 0x1E)
GOLD = RGBColor(0xC9, 0xA9, 0x6E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xB0, 0xC0)
MUTED = RGBColor(0x80, 0x80, 0x99)
TEAL = RGBColor(0x2D, 0xD4, 0xBF)
ROSE = RGBColor(0xF4, 0x4F, 0x9A)
PURPLE = RGBColor(0xA78B, 0xFA)[0:3] if False else RGBColor(0xA7, 0x8B, 0xFA)
BLUE = RGBColor(0x60, 0xA5, 0xFA)
AMBER = RGBColor(0xFB, 0xBF, 0x24)
GREEN = RGBColor(0x34, 0xD3, 0x99)


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=BG_CARD, corner_radius=Inches(0.15)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.rotation = 0.0
    # Corner radius via adjustments
    if shape.adjustments:
        try:
            shape.adjustments[0] = 0.05
        except Exception:
            pass
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, lines, default_size=16, default_color=LIGHT_GRAY):
    """lines = list of (text, size, color, bold)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        text, size, color, bold = line
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Segoe UI"
        p.space_after = Pt(4)
    return txBox


def add_icon_stat(slide, left, top, icon, value, label, accent=GOLD):
    """Add a stat card with icon, big number, and label"""
    card = add_shape(slide, left, top, Inches(2.6), Inches(1.5))
    add_text(slide, left + Inches(0.25), top + Inches(0.15), Inches(2.2), Inches(0.4), icon, font_size=22, color=accent)
    add_text(slide, left + Inches(0.25), top + Inches(0.5), Inches(2.2), Inches(0.5), value, font_size=28, color=WHITE, bold=True)
    add_text(slide, left + Inches(0.25), top + Inches(1.0), Inches(2.2), Inches(0.4), label, font_size=12, color=MUTED)


# ═══════════════════════════════════════════════
# SLIDE 1 — Team & Track
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

# Gold accent line
line_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(0.08), Inches(1.2))
line_shape.fill.solid()
line_shape.fill.fore_color.rgb = GOLD
line_shape.line.fill.background()

add_text(slide, Inches(1.2), Inches(0.6), Inches(8), Inches(0.6), "NATURALS BEAUTYTECH HACKATHON 2026", font_size=14, color=GOLD, bold=True)
add_text(slide, Inches(1.2), Inches(1.2), Inches(10), Inches(0.8), "AURA", font_size=60, color=WHITE, bold=True)
add_text(slide, Inches(1.2), Inches(2.0), Inches(10), Inches(0.5), "AI-Unified Response Architecture for Salon Intelligence", font_size=22, color=LIGHT_GRAY)

# Team card
add_shape(slide, Inches(1.2), Inches(3.2), Inches(5), Inches(3.2))
add_multi_text(slide, Inches(1.5), Inches(3.4), Inches(4.5), Inches(3.0), [
    ("TEAM", 12, GOLD, True),
    ("", 6, MUTED, False),
    ("Dhilip S E  |  Full-Stack AI Engineer", 16, WHITE, True),
    ("", 6, MUTED, False),
    ("Solo build: Backend + Frontend + AI Agents + DevOps", 13, LIGHT_GRAY, False),
    ("Stack: FastAPI + React + PostgreSQL + Gemini AI", 13, LIGHT_GRAY, False),
    ("Deployed: Vercel + Supabase", 13, LIGHT_GRAY, False),
])

# Track coverage card
add_shape(slide, Inches(6.8), Inches(3.2), Inches(5.5), Inches(3.2))
add_multi_text(slide, Inches(7.1), Inches(3.4), Inches(5.0), Inches(3.0), [
    ("ALL 6 TRACKS SOLVED", 12, GOLD, True),
    ("", 6, MUTED, False),
    ("Track 1  Service Standardization & Quality", 13, TEAL, False),
    ("Track 2  Staff Dependency & Development", 13, BLUE, False),
    ("Track 3  Hyper-Personalized Beauty", 13, ROSE, False),
    ("Track 4  AI Trend Prediction & Inventory", 13, AMBER, False),
    ("Track 5  Customer Experience Optimization", 13, GREEN, False),
    ("Track 6  Salon Business Intelligence", 13, PURPLE, False),
])

# Bottom stats
add_icon_stat(slide, Inches(1.2), Inches(6.6), "", "63", "AI Agents", TEAL)
add_icon_stat(slide, Inches(4.0), Inches(6.6), "", "60+", "Problems Solved", ROSE)
add_icon_stat(slide, Inches(6.8), Inches(6.6), "", "8", "Intelligence Modules", BLUE)
add_icon_stat(slide, Inches(9.6), Inches(6.6), "", "100+", "API Endpoints", GOLD)


# ═══════════════════════════════════════════════
# SLIDE 2 — The Problem
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.4), "THE PROBLEM", font_size=14, color=GOLD, bold=True)
add_text(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.7), "750 Salons. Zero Connected Intelligence.", font_size=40, color=WHITE, bold=True)
add_text(slide, Inches(0.8), Inches(1.5), Inches(10), Inches(0.5), "Every salon operates as an isolated island \u2014 no shared data, no AI, no standardization.", font_size=16, color=MUTED)

# Problem cards - Row 1
problems_row1 = [
    ("Inconsistent Quality", "Same service, different results at\nevery branch. No SOP tracking,\nno quality measurement.", ROSE),
    ("Staff = Single Point of Failure", "When a stylist leaves, customers\nfollow. Zero knowledge retention.\nNo transition planning.", BLUE),
    ("Generic Experiences", "Every customer gets the same\ntreatment. No personalization for\nhair type, skin, mood, or climate.", TEAL),
]

for i, (title, desc, accent) in enumerate(problems_row1):
    x = Inches(0.8) + Inches(i * 4.1)
    add_shape(slide, x, Inches(2.3), Inches(3.8), Inches(2.2))
    add_text(slide, x + Inches(0.3), Inches(2.45), Inches(3.2), Inches(0.15), "\u25cf", font_size=20, color=accent)
    add_text(slide, x + Inches(0.3), Inches(2.75), Inches(3.3), Inches(0.4), title, font_size=17, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.3), Inches(3.3), Inches(3.3), Inches(1.1), desc, font_size=13, color=LIGHT_GRAY)

# Problem cards - Row 2
problems_row2 = [
    ("Blind to Trends", "Trends emerge on social media.\nSalons discover them months late.\nInventory misaligned, staff untrained.", AMBER),
    ("Broken Customer Journey", "No follow-up after service. Long\nwaits with no engagement. No-shows\nwith no recovery. Language barriers.", GREEN),
    ("No Business Intelligence", "750 locations, zero visibility.\nCan't benchmark, can't predict\nchurn, can't measure training ROI.", PURPLE),
]

for i, (title, desc, accent) in enumerate(problems_row2):
    x = Inches(0.8) + Inches(i * 4.1)
    add_shape(slide, x, Inches(4.8), Inches(3.8), Inches(2.2))
    add_text(slide, x + Inches(0.3), Inches(4.95), Inches(3.2), Inches(0.15), "\u25cf", font_size=20, color=accent)
    add_text(slide, x + Inches(0.3), Inches(5.25), Inches(3.3), Inches(0.4), title, font_size=17, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.3), Inches(5.8), Inches(3.3), Inches(1.1), desc, font_size=13, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════
# SLIDE 3 — Root Cause Analysis
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.4), "ROOT CAUSE ANALYSIS", font_size=14, color=GOLD, bold=True)
add_text(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.7), "Why the Same Mistakes Repeat Every Month", font_size=40, color=WHITE, bold=True)

# Central flow diagram as cards with arrows
causes = [
    ("1. DATA SILOS", "Each salon uses paper/Excel.\nNo shared database across\n750 locations.", ROSE, "Customer history, preferences,\nallergies lost between visits"),
    ("2. NO AI LAYER", "Zero intelligent automation.\nEvery decision is manual,\ninconsistent, slow.", BLUE, "Stylists rely on memory.\nManagers rely on gut feeling."),
    ("3. NO STANDARDIZATION", "SOPs exist on paper but\nno tracking, no compliance\nmeasurement, no alerts.", TEAL, "Same haircut = different quality\nat different branches"),
    ("4. REACTIVE OPERATIONS", "Problems discovered after\ncustomer complaints. No\nprediction, no prevention.", AMBER, "Trends missed, staff churn\nsurprises, inventory stockouts"),
]

for i, (title, desc, accent, impact) in enumerate(causes):
    x = Inches(0.6) + Inches(i * 3.15)
    # Main cause card
    add_shape(slide, x, Inches(2.0), Inches(2.95), Inches(2.2))
    add_text(slide, x + Inches(0.2), Inches(2.15), Inches(2.55), Inches(0.4), title, font_size=15, color=accent, bold=True)
    add_text(slide, x + Inches(0.2), Inches(2.65), Inches(2.55), Inches(1.3), desc, font_size=13, color=LIGHT_GRAY)

    # Arrow down
    add_text(slide, x + Inches(1.1), Inches(4.25), Inches(0.8), Inches(0.4), "\u25bc", font_size=20, color=accent)

    # Impact card
    add_shape(slide, x, Inches(4.7), Inches(2.95), Inches(1.6), fill_color=RGBColor(0x1E, 0x1E, 0x2A))
    add_text(slide, x + Inches(0.2), Inches(4.8), Inches(2.55), Inches(0.3), "IMPACT", font_size=10, color=accent, bold=True)
    add_text(slide, x + Inches(0.2), Inches(5.15), Inches(2.55), Inches(1.0), impact, font_size=12, color=MUTED)

# Bottom conclusion
add_shape(slide, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.7), fill_color=RGBColor(0x1A, 0x15, 0x10))
add_text(slide, Inches(1.0), Inches(6.7), Inches(11), Inches(0.5),
         "RESULT:  Naturals loses customers, revenue, and competitive edge every single day these problems go unsolved.",
         font_size=15, color=GOLD, bold=True)


# ═══════════════════════════════════════════════
# SLIDE 4 — Our Solution: AURA
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.4), "OUR SOLUTION", font_size=14, color=GOLD, bold=True)
add_text(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.7), "AURA \u2014 One Platform. 8 AI Modules. 60 Problems Solved.", font_size=36, color=WHITE, bold=True)
add_text(slide, Inches(0.8), Inches(1.45), Inches(10), Inches(0.5), "AI-Unified Response Architecture  |  63 Intelligent Agents  |  Real-time Intelligence for Every Role", font_size=14, color=MUTED)

modules = [
    ("\U0001f4cb Beauty Passport", "Lifelong digital beauty identity.\nHair + skin diagnostics, allergies,\nlifestyle, goals, service history.", TEAL, "10 agents"),
    ("\U0001f9e0 SOULSKIN Engine", "World's first emotion-to-beauty\nsystem. 5 archetypes. Poetic soul\nreadings. Sensory design.", ROSE, "Integrated"),
    ("\U0001fa9e AR Smart Mirror", "Virtual try-on for hairstyles,\ncolors, makeup. MediaPipe +\nWebGL real-time rendering.", BLUE, "3D ready"),
    ("\U0001f3af AI Stylist Assistant", "Real-time SOP guidance. Step-by-\nstep coaching. Chemical safety\ngates. Rush detection.", GREEN, "12 agents"),
    ("\U0001f4c8 Trend Intelligence", "6-week early detection. Celebrity\nradar. Regional mapping. Inventory\nprediction. Campaign recs.", AMBER, "11 agents"),
    ("\U0001f465 Staff Intelligence", "Attrition prediction. Skill mapping.\nCareer paths. Knowledge capture.\nClient transition planning.", PURPLE, "11 agents"),
    ("\u2728 Experience Engine", "Smart waiting. Follow-ups. Real-time\ntransparency. No-show recovery.\nAmbient intelligence.", TEAL, "10 agents"),
    ("\U0001f4ca Salon BI Dashboard", "Franchise-wide analytics. LTV\nmodeling. Branch health prediction.\nTraining ROI. Portfolio optimizer.", GOLD, "10 agents"),
]

for i, (title, desc, accent, badge) in enumerate(modules):
    col = i % 4
    row = i // 4
    x = Inches(0.6) + Inches(col * 3.15)
    y = Inches(2.1) + Inches(row * 2.7)

    add_shape(slide, x, y, Inches(2.95), Inches(2.4))
    add_text(slide, x + Inches(0.2), y + Inches(0.15), Inches(2.55), Inches(0.4), title, font_size=14, color=accent, bold=True)
    add_text(slide, x + Inches(0.2), y + Inches(0.6), Inches(2.55), Inches(1.3), desc, font_size=12, color=LIGHT_GRAY)
    # Badge
    badge_shape = add_shape(slide, x + Inches(0.2), y + Inches(1.95), Inches(1.2), Inches(0.3), fill_color=RGBColor(0x25, 0x25, 0x30))
    add_text(slide, x + Inches(0.3), y + Inches(1.95), Inches(1.0), Inches(0.3), badge, font_size=10, color=accent, bold=True)


# ═══════════════════════════════════════════════
# SLIDE 5 — Implementation Model
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.4), "IMPLEMENTATION", font_size=14, color=GOLD, bold=True)
add_text(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.7), "From Concept to Pilot in 4 Weeks", font_size=40, color=WHITE, bold=True)

# Architecture overview
add_shape(slide, Inches(0.6), Inches(1.8), Inches(5.8), Inches(5.2))
add_text(slide, Inches(0.9), Inches(1.95), Inches(5.2), Inches(0.4), "TECH ARCHITECTURE", font_size=13, color=GOLD, bold=True)

arch_lines = [
    ("FRONTEND", 12, TEAL, True),
    ("React 19 + TypeScript + Vite 8 + PWA", 12, LIGHT_GRAY, False),
    ("TanStack Query + Zustand + Framer Motion", 12, LIGHT_GRAY, False),
    ("Recharts + Three.js + MediaPipe AR", 12, LIGHT_GRAY, False),
    ("", 8, MUTED, False),
    ("BACKEND", 12, BLUE, True),
    ("FastAPI (100% async) + SQLAlchemy 2.0", 12, LIGHT_GRAY, False),
    ("63 AI Agents across 6 tracks (9,729 LOC)", 12, LIGHT_GRAY, False),
    ("JWT Auth + RBAC (6 roles) + Celery tasks", 12, LIGHT_GRAY, False),
    ("", 8, MUTED, False),
    ("DATA & AI", 12, ROSE, True),
    ("PostgreSQL (Supabase) + Redis cache", 12, LIGHT_GRAY, False),
    ("Google Gemini 2.0 Flash (SOULSKIN AI)", 12, LIGHT_GRAY, False),
    ("30 database models + ACID compliant", 12, LIGHT_GRAY, False),
    ("", 8, MUTED, False),
    ("DEPLOYMENT", 12, GREEN, True),
    ("Vercel (frontend + backend serverless)", 12, LIGHT_GRAY, False),
    ("Supabase (PostgreSQL + Storage + Auth)", 12, LIGHT_GRAY, False),
    ("PWA: installable, offline-ready, push notifs", 12, LIGHT_GRAY, False),
]
add_multi_text(slide, Inches(0.9), Inches(2.35), Inches(5.2), Inches(4.5), arch_lines)

# Right side - Implementation timeline
add_shape(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(5.2))
add_text(slide, Inches(7.1), Inches(1.95), Inches(5.2), Inches(0.4), "4-WEEK PILOT PLAN", font_size=13, color=GOLD, bold=True)

weeks = [
    ("WEEK 1: Foundation", "Database setup + Auth + Core models\nUser roles (6 types) + Location setup\nBeauty Passport + Customer profiles", TEAL),
    ("WEEK 2: Intelligence", "63 AI agents deployed across 6 tracks\nSOULSKIN engine + AR Mirror\nSOP tracking + Quality scoring", BLUE),
    ("WEEK 3: Experience", "Smart queue + Follow-ups + Feedback\nTrend detection + Inventory prediction\nBI dashboard + Analytics", AMBER),
    ("WEEK 4: Production", "Vercel deployment + Supabase cloud\nPWA optimization + Push notifications\nPilot salon onboarding + testing", GREEN),
]

for i, (title, desc, accent) in enumerate(weeks):
    y = Inches(2.5) + Inches(i * 1.15)
    add_text(slide, Inches(7.1), y, Inches(0.3), Inches(0.3), "\u25cf", font_size=14, color=accent)
    add_text(slide, Inches(7.5), y - Inches(0.02), Inches(4.8), Inches(0.3), title, font_size=13, color=accent, bold=True)
    add_text(slide, Inches(7.5), y + Inches(0.3), Inches(4.8), Inches(0.8), desc, font_size=11, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════
# SLIDE 6 — Business Impact
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.4), "BUSINESS IMPACT", font_size=14, color=GOLD, bold=True)
add_text(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.7), "Every Metric That Moves the Bottom Line", font_size=40, color=WHITE, bold=True)

impacts = [
    ("100%", "SOP Compliance\nTracking", "Real-time step-by-step\nguidance with deviation\nalerts + quality scoring", TEAL),
    ("+40%", "Repeat\nBookings", "Personalized beauty\npassport + journey plans\n+ smart follow-ups", ROSE),
    ("6 wks", "Early Trend\nDetection", "Social + celebrity radar\nwith regional mapping\n+ inventory prediction", AMBER),
    ("80%", "Client Retention\non Staff Exit", "AI transition plans match\ncustomers to replacement\nstylists by skill + history", BLUE),
    ("+35 pts", "NPS Score\nImprovement", "Smart waiting + real-time\ntransparency + micro-\nfeedback + ambient AI", GREEN),
    ("30%", "Margin\nImprovement", "BI dashboard + LTV\nmodeling + portfolio\noptimizer + training ROI", PURPLE),
    ("50%", "Staff Retention\nBoost", "Career path visibility +\nattrition risk alerts +\nknowledge capture", GOLD),
    ("$500K+", "Annual Savings\nIdentified", "Inventory optimization +\nwaste reduction + smart\nscheduling + upsells", WHITE),
]

for i, (number, label, desc, accent) in enumerate(impacts):
    col = i % 4
    row = i // 4
    x = Inches(0.6) + Inches(col * 3.15)
    y = Inches(1.7) + Inches(row * 2.85)

    add_shape(slide, x, y, Inches(2.95), Inches(2.6))
    add_text(slide, x + Inches(0.25), y + Inches(0.2), Inches(2.5), Inches(0.6), number, font_size=36, color=accent, bold=True)
    add_text(slide, x + Inches(0.25), y + Inches(0.85), Inches(2.5), Inches(0.5), label, font_size=14, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.25), y + Inches(1.5), Inches(2.5), Inches(1.0), desc, font_size=11, color=MUTED)


# ═══════════════════════════════════════════════
# SLIDE 7 — Scalability & Data Strategy
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.4), "SCALABILITY & DATA", font_size=14, color=GOLD, bold=True)
add_text(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.7), "Built for 1 Salon Today. Ready for 1,000 Tomorrow.", font_size=40, color=WHITE, bold=True)

# Left - Scalability
add_shape(slide, Inches(0.6), Inches(1.8), Inches(5.8), Inches(5.2))
add_text(slide, Inches(0.9), Inches(1.95), Inches(5.2), Inches(0.4), "SCALABILITY ARCHITECTURE", font_size=13, color=GOLD, bold=True)

scale_lines = [
    ("SERVERLESS COMPUTE", 12, TEAL, True),
    ("Vercel Edge Functions \u2014 auto-scales to demand", 12, LIGHT_GRAY, False),
    ("Zero server management, pay-per-invocation", 12, LIGHT_GRAY, False),
    ("", 8, MUTED, False),
    ("DATABASE", 12, BLUE, True),
    ("Supabase PostgreSQL with connection pooling", 12, LIGHT_GRAY, False),
    ("30 models with ACID compliance + row-level security", 12, LIGHT_GRAY, False),
    ("Read replicas ready for 750+ concurrent locations", 12, LIGHT_GRAY, False),
    ("", 8, MUTED, False),
    ("CACHING & PERFORMANCE", 12, AMBER, True),
    ("Redis for sessions, queues, rate limiting", 12, LIGHT_GRAY, False),
    ("PWA with service workers \u2014 offline-first", 12, LIGHT_GRAY, False),
    ("CDN-served static assets via Vercel Edge", 12, LIGHT_GRAY, False),
    ("", 8, MUTED, False),
    ("MULTI-TENANCY", 12, GREEN, True),
    ("Location-scoped data isolation", 12, LIGHT_GRAY, False),
    ("Role-based access: 6 roles with granular permissions", 12, LIGHT_GRAY, False),
    ("Franchise \u2192 Regional \u2192 Location hierarchy", 12, LIGHT_GRAY, False),
]
add_multi_text(slide, Inches(0.9), Inches(2.35), Inches(5.2), Inches(4.5), scale_lines)

# Right - Data Strategy
add_shape(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(5.2))
add_text(slide, Inches(7.1), Inches(1.95), Inches(5.2), Inches(0.4), "DATA STRATEGY", font_size=13, color=GOLD, bold=True)

data_lines = [
    ("30 INTERCONNECTED MODELS", 12, TEAL, True),
    ("Complete salon ecosystem in one unified schema", 12, LIGHT_GRAY, False),
    ("Customer \u2192 Booking \u2192 Session \u2192 Quality \u2192 Feedback", 12, LIGHT_GRAY, False),
    ("", 8, MUTED, False),
    ("AI DATA PIPELINE", 12, ROSE, True),
    ("Every interaction feeds 63 AI agents", 12, LIGHT_GRAY, False),
    ("Trend signals from social + bookings + weather", 12, LIGHT_GRAY, False),
    ("Predictive models improve with every service", 12, LIGHT_GRAY, False),
    ("", 8, MUTED, False),
    ("PRIVACY & COMPLIANCE", 12, PURPLE, True),
    ("GDPR / India DPDP Act compliant", 12, LIGHT_GRAY, False),
    ("UUID keys prevent enumeration attacks", 12, LIGHT_GRAY, False),
    ("Bcrypt hashing + JWT + KMS encryption", 12, LIGHT_GRAY, False),
    ("", 8, MUTED, False),
    ("REAL-TIME INTELLIGENCE", 12, AMBER, True),
    ("Live queue, transparency feed, micro-feedback", 12, LIGHT_GRAY, False),
    ("Climate-aware recommendations (weather API)", 12, LIGHT_GRAY, False),
    ("6 languages: EN, Tamil, Hindi, Telugu, Kannada, ML", 12, LIGHT_GRAY, False),
]
add_multi_text(slide, Inches(7.1), Inches(2.35), Inches(5.2), Inches(4.5), data_lines)


# ═══════════════════════════════════════════════
# SLIDE 8 — Roadmap & Pilot Plan
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.4), "ROADMAP", font_size=14, color=GOLD, bold=True)
add_text(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.7), "6 Months from Prototype to Network Rollout", font_size=40, color=WHITE, bold=True)

# Timeline phases
phases = [
    ("MONTH 1-2", "PILOT", "Single salon deployment\n3-5 stylists onboarded\nBeauty Passport + SOULSKIN\nSOP tracking + Quality scoring\nCustomer feedback collection", TEAL, "NOW"),
    ("MONTH 3", "VALIDATE", "Measure all 8 modules\nA/B test vs. control salon\nStaff training effectiveness\nCustomer satisfaction delta\nBI dashboard calibration", BLUE, "Q2"),
    ("MONTH 4", "EXPAND", "5 salons across 2 cities\nRegional trend mapping live\nCross-location benchmarking\nFranchise BI dashboard\nAttrition prediction active", AMBER, "Q2"),
    ("MONTH 5-6", "SCALE", "50+ salons network-wide\nFull trend intelligence\nAR Mirror in flagship stores\nDigital Beauty Twin launch\nMobile app release", GREEN, "Q3"),
]

for i, (period, title, desc, accent, quarter) in enumerate(phases):
    x = Inches(0.4) + Inches(i * 3.2)

    # Phase card
    add_shape(slide, x, Inches(1.8), Inches(3.0), Inches(4.2))

    # Quarter badge
    badge = add_shape(slide, x + Inches(2.0), Inches(1.9), Inches(0.8), Inches(0.35), fill_color=RGBColor(0x25, 0x25, 0x30))
    add_text(slide, x + Inches(2.05), Inches(1.9), Inches(0.7), Inches(0.35), quarter, font_size=10, color=accent, bold=True, alignment=PP_ALIGN.CENTER)

    add_text(slide, x + Inches(0.25), Inches(1.95), Inches(2.5), Inches(0.3), period, font_size=11, color=MUTED)
    add_text(slide, x + Inches(0.25), Inches(2.3), Inches(2.5), Inches(0.5), title, font_size=24, color=accent, bold=True)
    add_text(slide, x + Inches(0.25), Inches(2.9), Inches(2.5), Inches(3.0), desc, font_size=12, color=LIGHT_GRAY)

# Bottom CTA
add_shape(slide, Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.9), fill_color=RGBColor(0x1A, 0x15, 0x10))
add_multi_text(slide, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.8), [
    ("AURA is not a concept \u2014 it's built, deployed, and ready for pilot.", 18, GOLD, True),
    ("63 agents  \u00b7  8 modules  \u00b7  60+ problems  \u00b7  All 6 tracks  \u00b7  One unified platform  \u00b7  Built by one engineer.", 13, LIGHT_GRAY, False),
])

# ── Save ──
output_path = r"c:\github\natural\AURA_Pitch_Deck.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
